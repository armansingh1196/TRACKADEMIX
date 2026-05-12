from __future__ import annotations

import json
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION_START
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt


BASE_DIR = Path(__file__).resolve().parents[1]
DELIVERABLES_DIR = BASE_DIR / "deliverables"
ARTIFACTS_DIR = BASE_DIR / "artifacts"
REPORT_DOCX = DELIVERABLES_DIR / "Student_Performance_AI_Report.docx"
PPTX_FILE = DELIVERABLES_DIR / "Student_Performance_AI_Presentation.pptx"
MODEL_SUMMARY_PATH = ARTIFACTS_DIR / "model_summary.json"
METRICS_PATH = ARTIFACTS_DIR / "metrics.json"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def add_bullets(document: Document, items: list[str]) -> None:
    for item in items:
        document.add_paragraph(item, style="List Bullet")


def build_report(summary: dict, metrics: dict) -> None:
    document = Document()
    style = document.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    title = document.add_paragraph()
    title.alignment = 1
    run = title.add_run("Student Performance AI\nReport-Ready Project Document")
    run.bold = True
    run.font.size = Pt(18)

    subtitle = document.add_paragraph()
    subtitle.alignment = 1
    subtitle.add_run("BIT Management System Prototype").italic = True

    document.add_paragraph("")
    document.add_heading("Abstract", level=1)
    document.add_paragraph(
        "This report presents a Student Performance AI prototype developed for a BIT Management System. "
        "The system predicts student performance bands, identifies at-risk students, and produces practical "
        "recommendations using a Random Forest model trained on a realistic synthetic dataset."
    )

    document.add_heading("1. Introduction", level=1)
    document.add_paragraph(
        "The objective of this project is to extend a student management system with AI-driven analytics "
        "and intervention support. Instead of only storing records, the system uses student data to estimate "
        "performance and highlight areas where academic support may be required."
    )

    document.add_heading("2. Problem Statement", level=1)
    document.add_paragraph(summary["problem_statement"])

    document.add_heading("3. Objectives", level=1)
    add_bullets(
        document,
        [
            "Predict student performance as Low, Medium, or High.",
            "Identify academically at-risk students.",
            "Generate recommendation-oriented intervention suggestions.",
            "Present outputs through a dashboard suitable for project review.",
            "Prepare the project for future system integration.",
        ],
    )

    document.add_heading("4. Dataset And Features", level=1)
    document.add_paragraph(
        "The prototype uses a realistic synthetic dataset because live institutional data was not connected "
        "during the prototype phase. The feature set includes academic, behavioral, and contextual variables."
    )
    document.add_paragraph("Feature columns used by the model:")
    add_bullets(document, summary["feature_columns"])

    document.add_heading("5. Model Choice: Random Forest", level=1)
    document.add_paragraph(
        "Random Forest was selected because it is reliable on structured tabular data, easier to explain in "
        "an academic review, robust against overfitting compared with a single decision tree, and capable of "
        "providing feature importance."
    )
    add_bullets(
        document,
        [
            "Works well on structured educational data.",
            "Captures non-linear relationships between student indicators.",
            "Improves stability by combining multiple trees.",
            "Provides feature importance for explainability.",
            "Gives strong baseline performance for a prototype.",
        ],
    )

    document.add_heading("6. How The Model Works", level=1)
    add_bullets(
        document,
        [
            "Student data is collected through input features.",
            "Numeric fields are imputed with median values.",
            "Categorical fields are imputed with the most frequent values.",
            "Categorical features are encoded using one-hot encoding.",
            "A Random Forest classifier predicts the performance band.",
            "Predictions are translated into at-risk information and recommendations.",
        ],
    )

    document.add_heading("7. Target Definition", level=1)
    for label, text in summary["target_definition"].items():
        document.add_paragraph(f"{label}: {text}", style="List Bullet")

    document.add_heading("8. Model Performance", level=1)
    document.add_paragraph(f"Accuracy: {metrics['accuracy']:.4f}")
    document.add_paragraph(f"Weighted F1-score: {metrics['weighted_f1']:.4f}")
    document.add_paragraph(f"Training rows: {metrics['train_rows']}")
    document.add_paragraph(f"Test rows: {metrics['test_rows']}")

    document.add_paragraph("Class-wise performance:")
    report = metrics["classification_report"]
    for label in ["Low", "Medium", "High"]:
        document.add_paragraph(
            f"{label}: precision={report[label]['precision']:.4f}, "
            f"recall={report[label]['recall']:.4f}, "
            f"f1-score={report[label]['f1-score']:.4f}",
            style="List Bullet",
        )

    document.add_heading("9. Confusion Matrix Interpretation", level=1)
    matrix = metrics["confusion_matrix"]
    table = document.add_table(rows=4, cols=4)
    table.style = "Table Grid"
    headers = ["Actual \\ Predicted", "Low", "Medium", "High"]
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    labels = metrics["class_labels"]
    for row_index, label in enumerate(labels, start=1):
        table.cell(row_index, 0).text = label
        for col_index, value in enumerate(matrix[row_index - 1], start=1):
            table.cell(row_index, col_index).text = str(value)

    document.add_paragraph(
        "The model predicts the Medium class most accurately, while some Low and High students are still "
        "being grouped into Medium. This is acceptable for a prototype and indicates a realistic classification challenge."
    )

    document.add_heading("10. Top Important Features", level=1)
    top_features = summary["top_features"][:5]
    for item in top_features:
        document.add_paragraph(
            f"{item['feature']}: importance={item['importance']:.4f}",
            style="List Bullet",
        )

    document.add_heading("11. Recommendation Layer", level=1)
    document.add_paragraph(
        "The model is combined with a rule-based recommendation engine. After prediction, the system checks "
        "attendance, exam performance, assignment quality, quiz scores, study hours, LMS activity, connectivity, "
        "and workload indicators to produce intervention-oriented recommendations."
    )
    add_bullets(
        document,
        [
            "Attendance mentoring",
            "Remedial internal-exam support",
            "Assignment tracking",
            "Topic-wise quiz practice",
            "Study planning and LMS engagement support",
        ],
    )

    document.add_heading("12. Dashboard And Demo Value", level=1)
    document.add_paragraph(
        "The dashboard displays model metrics, class distributions, feature importance, sample predictions, "
        "and a live student prediction form. This makes the prototype demonstration-ready."
    )
    for image_name in [
        "target_distribution.png",
        "feature_importance.png",
        "risk_distribution.png",
    ]:
        image_path = ARTIFACTS_DIR / image_name
        if image_path.exists():
            document.add_picture(str(image_path), width=Inches(5.8))

    document.add_heading("13. How To Run The Application", level=1)
    add_bullets(
        document,
        [
            "python -m venv .venv",
            ".venv\\Scripts\\activate",
            "pip install -r requirements.txt",
            "python -m src.student_performance_ai.data_generation",
            "python -m src.student_performance_ai.training",
            "streamlit run app/dashboard.py",
        ],
    )
    document.add_paragraph("You can also use the one-click launchers: run_app.bat or run_app.ps1.")

    document.add_heading("14. Limitations", level=1)
    add_bullets(
        document,
        [
            "The dataset is synthetic rather than live institutional data.",
            "Recommendations are rule-based and not directly learned by the model.",
            "The project is still a prototype and not yet integrated into a production backend.",
            "The current accuracy is suitable for review and demo, but not final deployment.",
        ],
    )

    document.add_heading("15. Conclusion", level=1)
    document.add_paragraph(
        "The Student Performance AI prototype demonstrates how a BIT Management System can be enhanced with "
        "AI-driven analytics and recommendation support. Random Forest is an appropriate model for this stage "
        "because it is explainable, effective on tabular data, and useful for demonstrating feature-driven insights."
    )

    document.save(REPORT_DOCX)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptPt(22)


def add_image_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), PptInches(0.8), PptInches(1.5), width=PptInches(8.0))
    tx_box = slide.shapes.add_textbox(PptInches(0.9), PptInches(6.7), PptInches(8.2), PptInches(0.6))
    tx_box.text_frame.text = caption


def build_presentation(summary: dict, metrics: dict) -> None:
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    add_title_slide(
        prs,
        "Student Performance AI",
        "BIT Management System Prototype",
    )
    add_bullet_slide(
        prs,
        "Project Objective",
        [
            "Predict student performance bands using academic and behavioral data.",
            "Identify at-risk students for early intervention.",
            "Generate recommendation-oriented support actions.",
        ],
    )
    add_bullet_slide(
        prs,
        "Why Random Forest",
        [
            "Strong fit for structured tabular data.",
            "More explainable than deep learning for review presentations.",
            "Captures non-linear relationships between student indicators.",
            "Provides feature importance for interpretation.",
        ],
    )
    add_bullet_slide(
        prs,
        "Input Features",
        [
            "Attendance, assignments, quizzes, internal exams, lab performance.",
            "Study hours, LMS logins, previous GPA, project score.",
            "Department, study mode, internet access, part-time job, extracurricular activity.",
        ],
    )
    add_bullet_slide(
        prs,
        "Target And Output",
        [
            "Target variable: performance_band.",
            f"Low: {summary['target_definition']['Low']}",
            f"Medium: {summary['target_definition']['Medium']}",
            f"High: {summary['target_definition']['High']}",
            "The system also maps predictions into at-risk support output.",
        ],
    )
    add_bullet_slide(
        prs,
        "Model Pipeline",
        [
            "Impute numeric values with median.",
            "Impute categorical values with most frequent values.",
            "Apply one-hot encoding to categorical features.",
            "Train RandomForestClassifier on the processed feature set.",
            "Generate predictions, risk output, and recommendations.",
        ],
    )
    add_bullet_slide(
        prs,
        "Current Results",
        [
            f"Accuracy: {metrics['accuracy']:.4f}",
            f"Weighted F1-score: {metrics['weighted_f1']:.4f}",
            f"Training rows: {metrics['train_rows']}",
            f"Test rows: {metrics['test_rows']}",
            "Medium class is predicted most strongly in the current prototype.",
        ],
    )

    for image_name, title, caption in [
        ("target_distribution.png", "Performance Distribution", "Distribution of Low, Medium, and High performance bands."),
        ("feature_importance.png", "Feature Importance", "Random Forest highlights the strongest performance drivers."),
        ("risk_distribution.png", "Predicted Risk Distribution", "Dashboard-ready view of at-risk prediction output."),
    ]:
        image_path = ARTIFACTS_DIR / image_name
        if image_path.exists():
            add_image_slide(prs, title, image_path, caption)

    add_bullet_slide(
        prs,
        "Recommendation Logic",
        [
            "Prediction is combined with rule-based intervention support.",
            "Low attendance triggers attendance mentoring.",
            "Weak internal exams trigger remedial support.",
            "Low assignments or quizzes trigger targeted academic guidance.",
            "Low engagement triggers study planning and LMS support.",
        ],
    )
    add_bullet_slide(
        prs,
        "Dashboard Demo",
        [
            "Shows model metrics and class distributions.",
            "Displays top important predictors.",
            "Presents sample predictions from test data.",
            "Allows live student input for real-time prediction and recommendations.",
        ],
    )
    add_bullet_slide(
        prs,
        "How To Run",
        [
            "Use run_app.bat or run_app.ps1 for one-click launch.",
            "Or run python -m src.student_performance_ai.training and streamlit run app/dashboard.py.",
            "Open the local Streamlit URL in the browser.",
        ],
    )
    add_bullet_slide(
        prs,
        "Conclusion",
        [
            "Random Forest is an appropriate review-ready model for this prototype.",
            "The system combines analytics, prediction, and recommendations.",
            "The next step is integration with real student data and backend APIs.",
        ],
    )

    prs.save(PPTX_FILE)


def main() -> None:
    DELIVERABLES_DIR.mkdir(parents=True, exist_ok=True)
    summary = load_json(MODEL_SUMMARY_PATH)
    metrics = load_json(METRICS_PATH)
    build_report(summary, metrics)
    build_presentation(summary, metrics)
    print(f"Created {REPORT_DOCX}")
    print(f"Created {PPTX_FILE}")


if __name__ == "__main__":
    main()
